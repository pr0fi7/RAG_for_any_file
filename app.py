import sys
import os

# Add the root directory to the Python path
sys.path.append(os.path.abspath('/rag_nxt_lvl/RagOnTheFly/'))

import streamlit as st
from PyPDF2 import PdfReader
from langchain.embeddings import OpenAIEmbeddings
from langchain.vectorstores import FAISS #local one so everything will be deleted after the session
from langchain_text_splitters import CharacterTextSplitter
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain.chat_models import ChatOpenAI
import docx
from pptx import Presentation
import zipfile
import os
from htmlTemplates import css, bot_template, user_template
from link_parser import LinkParser
import json
from langchain.prompts import ChatPromptTemplate, SystemMessagePromptTemplate
import pandas as pd

link_parser = LinkParser()

def get_xls_text(xls_files):
    text = ""
    for xls_file in xls_files:
        xls = pd.ExcelFile(xls_file)
        for sheet in xls.sheet_names:
            df = pd.read_excel(xls, sheet)
            text += df.to_string()
    return text

def unzip_file(zip_file, extract_to):
    with zipfile.ZipFile(zip_file, 'r') as zip_ref:
        zip_ref.extractall(extract_to)

def get_pdf_text(pdf_files):
    text = ""
    for pdf in pdf_files:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def getText(filename):
    doc = docx.Document(filename)
    fullText = [para.text for para in doc.paragraphs]
    return '\n'.join(fullText)

def get_pptx_text(pptx_files):
    text = ""
    for pptx_file in pptx_files:
        prs = Presentation(pptx_file)
        fullText = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    fullText.append(shape.text)
        text += '\n'.join(fullText)
    return text

def get_text_from_url(url):
    try:
        return link_parser.extract_text_from_website(url)
    except Exception as e:
        st.error(f"An error occurred: {str(e)}. Please enter a valid URL.")
        return None

def get_text_chunks(text, chunk_size, chunk_overlap):
    splitter = CharacterTextSplitter(separator="\n", chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    text_chunks = splitter.split_text(text)
    return text_chunks

def get_vector_store(text_chunks):
    embeddings = OpenAIEmbeddings()
    vector_store = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vector_store

def get_conversation_chain(vectorstore, prompt):
    llm = ChatOpenAI(temperature=0, model_name="gpt-4o")
    memory = ConversationBufferMemory(memory_key='chat_history', return_messages=True,  output_key='answer')

    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm, 
        retriever=vectorstore.as_retriever(), 
        memory=memory, 
        return_source_documents=True,
        combine_docs_chain_kwargs={"prompt": prompt}
    )
    return conversation_chain

def handle_userinput(user_question):
    if st.session_state.conversation is None:
        st.error("Please upload a document or provide a URL before asking a question.")
        return

    response = st.session_state.conversation({'question': user_question})
    print(response)
    st.session_state.chat_history = response['chat_history']

    for i, message in list(enumerate(st.session_state.chat_history)):
        if not i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

def get_source_documents(user_question):
    response = st.session_state.conversation({'question': user_question})
    print(response)
    st.session_state.chat_history = response['chat_history']

    res_dict = {
        "answer_from_llm": response["answer"],   ### looks up result key from raw output
    }
    
    res_dict["source_documents"] = []    ### create an empty array for source documents key front result dict

    for each_source in response["source_documents"]:
        res_dict["source_documents"].append({
            "page_content": each_source.page_content,
        })
    
    pretty_sources = json.dumps(res_dict["source_documents"], indent=4, default=str)
    
    print(pretty_sources)
    st.write(pretty_sources)

def main():
    prompt = """
    You are an AI assistant for answering questions about the provided context. 
    Below You are given contextual information and a question,  provide a conversational answer.
    If you don't know the answer, just say 'I am very sorry, there is no information on that topic in the provided context. 
    {context}
    Question: {question}
    Answer:
    """
    raw_text = ""  # Initialize raw_text variable
    api_key = st.secrets["OPENAI_API_KEY"]
    st.set_page_config(page_title="RAG", page_icon=":shark:")
    st.write(css, unsafe_allow_html=True)
    vector_store = None

    if 'conversation' not in st.session_state or st.session_state.conversation is None:
        st.session_state.conversation = None
    st.header("RAG On the Fly with any document / URL :shark:")

    # Check if the 'user_question' key exists in the session state
    if 'user_question' not in st.session_state:
        # If not, create it and set its initial value to an empty string
        st.session_state.user_question = ""

    # Use the 'user_question' key to refer to the original text input field

    message = st.chat_input("Ask a question")
    if st.button("Get Source Documents"):
        get_source_documents(st.session_state.user_question)

    if message and message != st.session_state.user_question:
        handle_userinput(message)

    with st.sidebar:
        st.subheader("Load your documents")
        st.subheader('Current version, supports PDF, DOCX, PPTX, ZIP and URL')

        docs = st.file_uploader("Upload a document", accept_multiple_files=True)
        url = st.text_input("Enter a URL")

        chunk_size = st.slider("Chunk size", value=1200, min_value=200, max_value=2000, step=100)
        overlap = st.slider("Overlap", value=200, min_value=0, max_value=400, step=50)

        prompt_template = st.text_area("System Prompt Template")


        if st.button("RAG it now !"):
            with st.spinner("Analyzing, Vectorizing, Retrieving..."):
                if docs:
                    pdf_files = []
                    pptx_files = []
                    docx_files = []
                    xlsx_files = []

                    for doc in docs:
                        if doc.name.endswith('.zip'):
                            unzip_file(doc, 'temp')

                            extracted_files = []
                            for root, dirs, files in os.walk('temp'):
                                for file in files:
                                    extracted_files.append(os.path.join(root, file))

                            print(f"Extracted files: {extracted_files}")
                            for file in extracted_files:
                                print(f"Processing file: {file}")
                                if file.endswith('.pdf'):
                                    pdf_files.append(file)
                                elif file.endswith('.docx'):
                                    docx_files.append(file)
                                elif file.endswith('.pptx'):
                                    pptx_files.append(file)
                        elif doc.name.endswith('.pdf'):
                            pdf_files.append(doc)
                        elif doc.name.endswith('.docx'):
                            docx_files.append(doc)
                        elif doc.name.endswith('.pptx'):
                            pptx_files.append(doc)
                        elif doc.name.endswith('.xlsx'):
                            xlsx_files.append(doc)

                    if pdf_files:
                        raw_text += get_pdf_text(pdf_files)
                    if docx_files:
                        for docx_file in docx_files:
                            raw_text += getText(docx_file)
                    if pptx_files:
                        raw_text += get_pptx_text(pptx_files)
                    if xlsx_files:
                        raw_text += get_xls_text(xlsx_files)

                elif url:
                    raw_text += get_text_from_url(url)

                else:
                    st.error("Please upload a PDF, DOCX, PPTX, XLSX file or a ZIP file containing such file(s) or provide a URL for analysis.")
                    st.stop()

                if raw_text is None:
                    st.error("No text extracted from the documents.")
                    st.stop()
                

                text_chunks = get_text_chunks(raw_text, chunk_size, overlap)
                vectorstore = get_vector_store(text_chunks)

                if not prompt_template:
                    prompt_template = prompt
                else:
                    prompt_template = prompt_template + "context: {context} \nQuestion: {question} \nAnswer:"
                print(prompt_template)

                messages = [
                SystemMessagePromptTemplate.from_template(prompt_template),
                ]

                qa_prompt = ChatPromptTemplate.from_messages( messages )

                # putting conversation in session state to keep it alive for the session but it doesn't work
                st.session_state.conversation = get_conversation_chain(vectorstore, qa_prompt)
                st.write("Text extraction completed! \n Ask a question in the chatbox to get started.")

        if st.button('Clear Context'):
            # Reset the FAISS database / vector store
            vector_store = None
            st.session_state.conversation = None    
            st.session_state.user_question = ""
            st.session_state.chat_history = []
            conversation_chain = None
            st.write("Context has been cleared.")

        st.subheader("Manual")
        st.write("1. Upload a PDF, DOCX, or PPTX file or provide a URL for analysis.")
        st.write('Note: you can upload multiple files at once as well as combine different file types and content from URLs.')
        st.write("2. Click the 'Analyze' button to start the analysis.")
        st.write("3. Ask a question in the chatbox to get started.")
        st.write("4. The chatbot will provide answers based on the content of the uploaded document or URL.")

if __name__ == '__main__':
    main()
